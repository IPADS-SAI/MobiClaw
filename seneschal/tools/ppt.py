# -*- coding: utf-8 -*-
"""PowerPoint / PPTX tools for Seneschal.

Provides async tool functions for reading, creating, and editing PPTX files
via python-pptx (lazy-imported so the import error is handled gracefully).

Functions
---------
read_pptx_summary          – 读取 PPTX，返回每张幻灯片的标题/正文/备注/形状/图片信息
create_pptx_from_outline   – 从幻灯片大纲列表创建新 PPTX，支持模板与全局字体默认值
edit_pptx                  – 综合编辑：全局文本替换 / 追加幻灯片 / 按索引删除幻灯片
insert_pptx_image          – 向指定幻灯片插入图片，支持英寸定位与尺寸
set_pptx_text_style        – 在指定幻灯片搜索文本，对匹配 run 应用字体样式

Slide dict schema (used by create_pptx_from_outline / edit_pptx.add_slides)
---------------------------------------------------------------------------
{
  "title":      str,               # 幻灯片标题
  "content":    str | list[str],   # 正文/要点，列表时每项为一行
  "notes":      str,               # 演讲者备注
  "layout":     int,               # 布局索引（0=标题页，1=标题+内容，默认 1）
  "images": [
    {"path": str, "left": float, "top": float, "width": float, "height": float}
  ],                               # 英寸单位，width/height 可省略（保持原始比例）
  "font_size":  float,             # 正文字号 pt，覆盖全局默认
  "font_color": str,               # "#RRGGBB" 十六进制颜色，覆盖全局默认
  "bold":       bool,
  "italic":     bool,
}
"""

from __future__ import annotations

import os
import tempfile
import urllib.request
from pathlib import Path
from typing import Any

from agentscope.message import TextBlock
from agentscope.tool import ToolResponse


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _resolve_write_path(path: str) -> tuple[Path | None, str | None]:
    """Validate and resolve *output_path*, respecting SENESCHAL_FILE_WRITE_ROOT."""
    resolved_path = (path or "").strip()
    if not resolved_path:
        return None, "empty_path"

    root = (os.environ.get("SENESCHAL_FILE_WRITE_ROOT") or "").strip()
    target = Path(resolved_path).expanduser()
    if root:
        root_path = Path(root).expanduser().resolve()
        target = target if target.is_absolute() else root_path / target
        try:
            target = target.resolve()
        except FileNotFoundError:
            target = target.absolute()
        if target != root_path and root_path not in target.parents:
            return None, "path_outside_root"
    else:
        target = target.resolve()
    return target, None


def _parse_rgb_color(hex_color: str | None) -> Any | None:
    """Convert '#RRGGBB' string to pptx.dml.color.RGBColor, or None if invalid."""
    if not hex_color:
        return None
    try:
        from pptx.dml.color import RGBColor  # type: ignore
        cleaned = hex_color.lstrip("#")
        if len(cleaned) != 6:
            return None
        r = int(cleaned[0:2], 16)
        g = int(cleaned[2:4], 16)
        b = int(cleaned[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return None


def _slide_text_lines(slide: Any) -> list[str]:
    """Return non-empty text lines from every text frame in a slide."""
    lines: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    lines.append(line)
    return lines


def _is_url(value: str) -> bool:
    text = (value or "").strip().lower()
    return text.startswith("http://") or text.startswith("https://")


def _download_image_to_temp(url: str) -> tuple[Path | None, str | None]:
    """Download an image URL to a temporary file and return its path."""
    try:
        suffix = Path(url.split("?", 1)[0]).suffix or ".img"
        with urllib.request.urlopen(url, timeout=20) as resp:
            data = resp.read()
        if not data:
            return None, "empty_download"
        fd, tmp_path = tempfile.mkstemp(prefix="seneschal_ppt_", suffix=suffix)
        os.close(fd)
        target = Path(tmp_path)
        target.write_bytes(data)
        return target, None
    except Exception as exc:
        return None, str(exc)


def _apply_run_style(
    text_frame: Any,
    font_size: float | None,
    font_color: str | None,
    bold: bool | None,
    italic: bool | None,
) -> None:
    """Apply font style to every run in a text frame."""
    from pptx.util import Pt  # type: ignore

    color = _parse_rgb_color(font_color)
    for para in text_frame.paragraphs:
        for run in para.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            if color is not None:
                run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic


def _find_body_text_frame(slide: Any) -> Any | None:
    """Find the best candidate text frame for body content in current slide."""
    # Prefer placeholders that have a text frame and are not the title placeholder.
    for shape in slide.placeholders:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False) and shape == slide.shapes.title:
            continue
        return shape.text_frame

    # Fallback to any non-title text shape.
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape == slide.shapes.title:
            continue
        return shape.text_frame
    return None


def _write_content_to_text_frame(text_frame: Any, content: str | list[str]) -> None:
    """Write string/list content into a text frame with paragraph support."""
    if isinstance(content, list):
        lines = [str(item).strip() for item in content if str(item).strip()]
    else:
        text = str(content or "").strip()
        lines = [line.strip() for line in text.splitlines() if line.strip()]

    if not lines:
        return

    text_frame.clear()
    first_para = text_frame.paragraphs[0]
    first_para.text = lines[0]
    for line in lines[1:]:
        para = text_frame.add_paragraph()
        para.text = line


def _add_body_textbox(slide: Any, content: str | list[str]) -> Any:
    """Create a fallback textbox and write content into it."""
    from pptx.util import Inches  # type: ignore

    # A safe default area for most 16:9 and 4:3 templates.
    textbox = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(1.8),
        Inches(8.4),
        Inches(3.8),
    )
    text_frame = textbox.text_frame
    _write_content_to_text_frame(text_frame, content)
    return text_frame


def _resolve_image_path(image_input: str) -> tuple[Path | None, str | None, bool]:
    """Resolve local path or URL into a local image file path.

    Returns: (path, error, is_temp_download)
    """
    raw = (image_input or "").strip()
    if not raw:
        return None, "empty_image_path", False

    if _is_url(raw):
        downloaded, err = _download_image_to_temp(raw)
        if err:
            return None, f"download_failed:{err}", False
        return downloaded, None, True

    path = Path(raw).expanduser()
    if not path.exists() or not path.is_file():
        return None, "image_not_found", False
    return path, None, False


def _add_slide_from_dict(prs: Any, slide_data: dict[str, Any]) -> dict[str, Any]:
    """Append a new slide from dict and return structured diagnostics.

    Returns dict keys: slide, content_mode, warnings, image_results
    """
    from pptx.util import Inches  # type: ignore

    layout_index = int(slide_data.get("layout", 1))
    try:
        slide_layout = prs.slide_layouts[layout_index]
    except IndexError:
        slide_layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]

    slide = prs.slides.add_slide(slide_layout)
    warnings: list[str] = []
    image_results: list[dict[str, Any]] = []
    temp_files: list[Path] = []

    title_text = str(slide_data.get("title") or "").strip()
    content = slide_data.get("content") or ""
    has_content = bool(content)

    font_size = slide_data.get("font_size")
    font_color = slide_data.get("font_color")
    bold = slide_data.get("bold")
    italic = slide_data.get("italic")

    # --- title placeholder ---
    if slide.shapes.title and title_text:
        slide.shapes.title.text = title_text

    # --- body/content: try placeholder/text-frame first, fallback to textbox ---
    content_mode = "none"
    if has_content:
        body_text_frame = _find_body_text_frame(slide)
        if body_text_frame is not None:
            _write_content_to_text_frame(body_text_frame, content)
            content_mode = "placeholder"
            if any(v is not None for v in [font_size, font_color, bold, italic]):
                _apply_run_style(body_text_frame, font_size, font_color, bold, italic)
        else:
            warnings.append("content_placeholder_missing_fallback_textbox")
            textbox_tf = _add_body_textbox(slide, content)
            content_mode = "textbox_fallback"
            if any(v is not None for v in [font_size, font_color, bold, italic]):
                _apply_run_style(textbox_tf, font_size, font_color, bold, italic)

    # --- notes ---
    notes_text = str(slide_data.get("notes") or "").strip()
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

    # --- images ---
    for idx, img in enumerate(slide_data.get("images") or [], start=1):
        if not isinstance(img, dict):
            warnings.append(f"image_{idx}_invalid_schema")
            image_results.append({"index": idx, "status": "error", "reason": "invalid_schema"})
            continue

        source = str(img.get("path") or img.get("url") or "").strip()
        resolved_img, err, is_temp = _resolve_image_path(source)
        if err or resolved_img is None:
            reason = err or "image_resolve_failed"
            warnings.append(f"image_{idx}_{reason}")
            image_results.append(
                {
                    "index": idx,
                    "status": "error",
                    "reason": reason,
                    "source": source,
                }
            )
            continue

        if is_temp:
            temp_files.append(resolved_img)

        left = Inches(float(img.get("left", 1.0)))
        top = Inches(float(img.get("top", 1.0)))
        kwargs: dict[str, Any] = {"left": left, "top": top}
        if img.get("width") is not None:
            kwargs["width"] = Inches(float(img["width"]))
        if img.get("height") is not None:
            kwargs["height"] = Inches(float(img["height"]))

        try:
            slide.shapes.add_picture(str(resolved_img), **kwargs)
            image_results.append(
                {
                    "index": idx,
                    "status": "ok",
                    "source": source,
                    "resolved_path": str(resolved_img),
                }
            )
        except Exception as exc:
            warnings.append(f"image_{idx}_insert_failed")
            image_results.append(
                {
                    "index": idx,
                    "status": "error",
                    "reason": f"insert_failed:{exc}",
                    "source": source,
                }
            )

    return {
        "slide": slide,
        "content_mode": content_mode,
        "warnings": warnings,
        "image_results": image_results,
        "temp_files": temp_files,
    }


# ---------------------------------------------------------------------------
# Public tool functions
# ---------------------------------------------------------------------------


async def read_pptx_summary(
    file_path: str,
    max_slides: int | None = None,
) -> ToolResponse:
    """读取 PPTX/PPT 文件，返回每张幻灯片的标题、正文文本、备注、形状数量和图片数量。

    Parameters
    ----------
    file_path:  本地 PPTX/PPT 文件路径。
    max_slides: 最多读取的幻灯片数量，默认读取全部。
    """
    resolved_path = (file_path or "").strip()
    if not resolved_path:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] Empty file path.")],
            metadata={"error": "empty_path"},
        )

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] python-pptx is not installed. Run: pip install python-pptx")],
            metadata={"error": "missing_dependency"},
        )

    target = Path(resolved_path).expanduser()
    if not target.exists() or not target.is_file():
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] File not found: {target}")],
            metadata={"error": "not_found", "path": str(target)},
        )
    if target.suffix.lower() not in {".pptx", ".ppt"}:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Unsupported extension '{target.suffix}'. Expected .pptx or .ppt")],
            metadata={"error": "invalid_extension", "path": str(target)},
        )

    try:
        prs = Presentation(str(target))
        total_slides = len(prs.slides)
        limit = total_slides if max_slides is None else min(int(max_slides), total_slides)

        slide_summaries: list[dict[str, Any]] = []
        for idx, slide in enumerate(list(prs.slides)[:limit], start=1):
            title = (slide.shapes.title.text or "").strip() if slide.shapes.title else ""
            all_lines = _slide_text_lines(slide)
            body_lines = [ln for ln in all_lines if ln != title][:20]

            # MSO_SHAPE_TYPE.PICTURE == 13
            image_count = sum(1 for sh in slide.shapes if sh.shape_type == 13)

            notes_text = ""
            try:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                pass

            slide_summaries.append(
                {
                    "slide": idx,
                    "title": title,
                    "texts": body_lines,
                    "notes": notes_text[:200],
                    "shape_count": len(slide.shapes),
                    "image_count": image_count,
                }
            )

        lines = [f"[PPTX] {target.name}  total_slides={total_slides}  shown={limit}"]
        for s in slide_summaries:
            lines.append(f"  [{s['slide']}] title={s['title'] or '(none)'}  shapes={s['shape_count']}  images={s['image_count']}")
            for t in s["texts"]:
                lines.append(f"       • {t[:120]}")
            if s["notes"]:
                lines.append(f"       notes: {s['notes'][:80]}")

        return ToolResponse(
            content=[TextBlock(type="text", text="\n".join(lines))],
            metadata={
                "path": str(target),
                "total_slides": total_slides,
                "shown_slides": limit,
                "slides": slide_summaries,
            },
        )
    except Exception as exc:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Error reading file: {exc}")],
            metadata={"error": str(exc)},
        )


async def create_pptx_from_outline(
    output_path: str,
    slides: list[dict[str, Any]],
    template_path: str | None = None,
    default_font_size: float | None = None,
    default_font_color: str | None = None,
) -> ToolResponse:
    """从幻灯片大纲列表创建新 PPTX 文件，支持模板与全局字体默认值。

    Parameters
    ----------
    output_path:        输出文件路径，须以 .pptx 或 .ppt 结尾。
    slides:             幻灯片列表，每项为 slide dict（见模块文档）。
    template_path:      可选，PPTX 模板文件路径；不提供时使用库内置空白演示文稿。
    default_font_size:  全局默认正文字号（pt），各幻灯片可通过 "font_size" 字段覆盖。
    default_font_color: 全局默认正文颜色（"#RRGGBB"），各幻灯片可通过 "font_color" 字段覆盖。
    """
    target, error = _resolve_write_path(output_path)
    if error:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Path error: {error}")],
            metadata={"error": error},
        )
    if target.suffix.lower() not in {".pptx", ".ppt"}:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Output path must end in .pptx or .ppt: {target}")],
            metadata={"error": "invalid_extension", "path": str(target)},
        )

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] python-pptx is not installed. Run: pip install python-pptx")],
            metadata={"error": "missing_dependency"},
        )

    if not isinstance(slides, list) or not slides:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] 'slides' must be a non-empty list.")],
            metadata={"error": "empty_slides"},
        )

    for idx, slide_data in enumerate(slides, start=1):
        if not isinstance(slide_data, dict):
            return ToolResponse(
                content=[TextBlock(type="text", text=f"[PPTX] slides[{idx}] must be a dict.")],
                metadata={"error": "schema_invalid", "slide_index": idx},
            )
        has_title = bool(str(slide_data.get("title") or "").strip())
        has_content = bool(slide_data.get("content"))
        if not (has_title or has_content):
            return ToolResponse(
                content=[TextBlock(type="text", text=f"[PPTX] slides[{idx}] must have at least one of title/content.")],
                metadata={"error": "schema_invalid", "slide_index": idx},
            )

    try:
        if template_path:
            tmpl = Path(template_path).expanduser()
            if not tmpl.exists():
                return ToolResponse(
                    content=[TextBlock(type="text", text=f"[PPTX] Template not found: {tmpl}")],
                    metadata={"error": "template_not_found", "path": str(tmpl)},
                )
            prs = Presentation(str(tmpl))
        else:
            prs = Presentation()

        warnings: list[str] = []
        content_modes: list[dict[str, Any]] = []
        image_results: list[dict[str, Any]] = []
        temp_files: list[Path] = []

        for i, slide_data in enumerate(slides, start=1):
            merged: dict[str, Any] = dict(slide_data)
            if default_font_size is not None and "font_size" not in merged:
                merged["font_size"] = default_font_size
            if default_font_color is not None and "font_color" not in merged:
                merged["font_color"] = default_font_color
            result = _add_slide_from_dict(prs, merged)
            warnings.extend(result.get("warnings") or [])
            content_modes.append({"slide": i, "mode": result.get("content_mode")})
            for item in result.get("image_results") or []:
                image_results.append({"slide": i, **item})
            temp_files.extend(result.get("temp_files") or [])

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
        for temp_file in temp_files:
            try:
                temp_file.unlink(missing_ok=True)
            except OSError:
                pass
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Wrote: {target}  slides={len(slides)}  warnings={len(warnings)}")],
            metadata={
                "path": str(target),
                "slides": len(slides),
                "warnings": warnings,
                "content_modes": content_modes,
                "image_results": image_results,
            },
        )
    except Exception as exc:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Error creating file: {exc}")],
            metadata={"error": str(exc)},
        )


async def edit_pptx(
    file_path: str,
    output_path: str,
    replacements: list[dict[str, str]] | None = None,
    add_slides: list[dict[str, Any]] | None = None,
    delete_slide_indices: list[int] | None = None,
) -> ToolResponse:
    """综合编辑 PPTX：全局文本替换、追加幻灯片、按索引删除幻灯片。

    Parameters
    ----------
    file_path:             源 PPTX 文件路径。
    output_path:           输出文件路径（可与 file_path 相同以原地保存）。
    replacements:          [{"old": "...", "new": "..."}] 全局文本替换，跨所有幻灯片。
    add_slides:            要追加的幻灯片列表（同 create_pptx_from_outline 的 slide dict）。
    delete_slide_indices:  要删除的幻灯片 1-based 索引列表（从后往前删以保持其他索引稳定）。
    """
    resolved_src = (file_path or "").strip()
    if not resolved_src:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] Empty source file path.")],
            metadata={"error": "empty_path"},
        )
    target, error = _resolve_write_path(output_path)
    if error:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Output path error: {error}")],
            metadata={"error": error},
        )

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] python-pptx is not installed. Run: pip install python-pptx")],
            metadata={"error": "missing_dependency"},
        )

    src = Path(resolved_src).expanduser()
    if not src.exists() or not src.is_file():
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Source file not found: {src}")],
            metadata={"error": "not_found", "path": str(src)},
        )

    try:
        prs = Presentation(str(src))

        # 1) Global text replacements
        replace_count = 0
        for repl in replacements or []:
            old = str(repl.get("old") or "")
            new = str(repl.get("new") or "")
            if not old:
                continue
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                replace_count += 1

        # 2) Delete slides by 1-based index (descending to preserve indices)
        delete_count = 0
        if delete_slide_indices:
            sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
            sorted_del = sorted(
                [i for i in delete_slide_indices if 1 <= i <= len(prs.slides)],
                reverse=True,
            )
            for idx in sorted_del:
                sldIdLst.remove(sldIdLst[idx - 1])
                delete_count += 1

        # 3) Append new slides
        added_count = 0
        warnings: list[str] = []
        image_results: list[dict[str, Any]] = []
        content_modes: list[dict[str, Any]] = []
        temp_files: list[Path] = []
        for slide_data in add_slides or []:
            if isinstance(slide_data, dict):
                result = _add_slide_from_dict(prs, slide_data)
                warnings.extend(result.get("warnings") or [])
                content_modes.append({"slide": len(prs.slides), "mode": result.get("content_mode")})
                for item in result.get("image_results") or []:
                    image_results.append({"slide": len(prs.slides), **item})
                temp_files.extend(result.get("temp_files") or [])
                added_count += 1

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
        for temp_file in temp_files:
            try:
                temp_file.unlink(missing_ok=True)
            except OSError:
                pass

        summary = (
            f"[PPTX] Wrote: {target}  "
            f"replacements={replace_count}  deleted={delete_count}  added={added_count}  "
            f"final_slides={len(prs.slides)}"
        )
        return ToolResponse(
            content=[TextBlock(type="text", text=summary)],
            metadata={
                "path": str(target),
                "replacements": replace_count,
                "deleted": delete_count,
                "added": added_count,
                "final_slides": len(prs.slides),
                "warnings": warnings,
                "content_modes": content_modes,
                "image_results": image_results,
            },
        )
    except Exception as exc:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Error editing file: {exc}")],
            metadata={"error": str(exc)},
        )


async def insert_pptx_image(
    file_path: str,
    output_path: str,
    slide_index: int,
    image_path: str,
    left_inches: float = 1.0,
    top_inches: float = 1.0,
    width_inches: float | None = None,
    height_inches: float | None = None,
) -> ToolResponse:
    """向指定幻灯片插入图片，支持英寸单位定位与尺寸。

    Parameters
    ----------
    file_path:     源 PPTX 文件路径。
    output_path:   输出文件路径。
    slide_index:   1-based 幻灯片索引。
    image_path:    本地图片路径（支持 PNG/JPG/GIF/BMP 等常见格式）。
    left_inches:   图片左边距（英寸），默认 1.0。
    top_inches:    图片上边距（英寸），默认 1.0。
    width_inches:  图片宽度（英寸），省略时保持原始比例。
    height_inches: 图片高度（英寸），省略时保持原始比例。
    """
    resolved_src = (file_path or "").strip()
    if not resolved_src:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] Empty source file path.")],
            metadata={"error": "empty_path"},
        )
    target, error = _resolve_write_path(output_path)
    if error:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Output path error: {error}")],
            metadata={"error": error},
        )
    img_resolved = (image_path or "").strip()
    if not img_resolved:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] Empty image path.")],
            metadata={"error": "empty_image_path"},
        )

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] python-pptx is not installed. Run: pip install python-pptx")],
            metadata={"error": "missing_dependency"},
        )

    src = Path(resolved_src).expanduser()
    if not src.exists() or not src.is_file():
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Source file not found: {src}")],
            metadata={"error": "not_found", "path": str(src)},
        )
    img_file = Path(img_resolved).expanduser()
    if not img_file.exists() or not img_file.is_file():
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Image file not found: {img_file}")],
            metadata={"error": "image_not_found", "path": str(img_file)},
        )

    try:
        prs = Presentation(str(src))
        total = len(prs.slides)
        if not 1 <= slide_index <= total:
            return ToolResponse(
                content=[TextBlock(type="text", text=f"[PPTX] slide_index {slide_index} out of range (1–{total}).")],
                metadata={"error": "slide_index_out_of_range", "total": total},
            )

        slide = prs.slides[slide_index - 1]
        kwargs: dict[str, Any] = {
            "left": Inches(left_inches),
            "top": Inches(top_inches),
        }
        if width_inches is not None:
            kwargs["width"] = Inches(width_inches)
        if height_inches is not None:
            kwargs["height"] = Inches(height_inches)
        slide.shapes.add_picture(str(img_file), **kwargs)

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Wrote: {target}  image inserted at slide {slide_index}")],
            metadata={
                "path": str(target),
                "slide_index": slide_index,
                "image": str(img_file),
                "left_inches": left_inches,
                "top_inches": top_inches,
                "width_inches": width_inches,
                "height_inches": height_inches,
            },
        )
    except Exception as exc:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Error inserting image: {exc}")],
            metadata={"error": str(exc)},
        )


async def set_pptx_text_style(
    file_path: str,
    output_path: str,
    slide_index: int,
    search_text: str,
    font_size: float | None = None,
    font_color: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    underline: bool | None = None,
) -> ToolResponse:
    """在指定幻灯片中搜索文本，对所有包含该文本的 run 应用字体样式。

    Parameters
    ----------
    file_path:   源 PPTX 文件路径。
    output_path: 输出文件路径。
    slide_index: 1-based 幻灯片索引。
    search_text: 要匹配的文字子串。
    font_size:   字号（pt），省略时不修改。
    font_color:  十六进制颜色字符串 "#RRGGBB"，省略时不修改。
    bold:        True/False，省略时不修改。
    italic:      True/False，省略时不修改。
    underline:   True/False，省略时不修改。
    """
    resolved_src = (file_path or "").strip()
    if not resolved_src:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] Empty source file path.")],
            metadata={"error": "empty_path"},
        )
    target, error = _resolve_write_path(output_path)
    if error:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Output path error: {error}")],
            metadata={"error": error},
        )
    if not (search_text or "").strip():
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] search_text must not be empty.")],
            metadata={"error": "empty_search_text"},
        )

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore
    except ImportError:
        return ToolResponse(
            content=[TextBlock(type="text", text="[PPTX] python-pptx is not installed. Run: pip install python-pptx")],
            metadata={"error": "missing_dependency"},
        )

    src = Path(resolved_src).expanduser()
    if not src.exists() or not src.is_file():
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Source file not found: {src}")],
            metadata={"error": "not_found", "path": str(src)},
        )

    try:
        prs = Presentation(str(src))
        total = len(prs.slides)
        if not 1 <= slide_index <= total:
            return ToolResponse(
                content=[TextBlock(type="text", text=f"[PPTX] slide_index {slide_index} out of range (1–{total}).")],
                metadata={"error": "slide_index_out_of_range", "total": total},
            )

        slide = prs.slides[slide_index - 1]
        color = _parse_rgb_color(font_color)
        modified_runs = 0

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if search_text in run.text:
                        if font_size is not None:
                            run.font.size = Pt(font_size)
                        if color is not None:
                            run.font.color.rgb = color
                        if bold is not None:
                            run.font.bold = bold
                        if italic is not None:
                            run.font.italic = italic
                        if underline is not None:
                            run.font.underline = underline
                        modified_runs += 1

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
        return ToolResponse(
            content=[TextBlock(
                type="text",
                text=(
                    f"[PPTX] Wrote: {target}  "
                    f"styled {modified_runs} run(s) matching '{search_text}' on slide {slide_index}"
                ),
            )],
            metadata={
                "path": str(target),
                "slide_index": slide_index,
                "search_text": search_text,
                "modified_runs": modified_runs,
            },
        )
    except Exception as exc:
        return ToolResponse(
            content=[TextBlock(type="text", text=f"[PPTX] Error styling text: {exc}")],
            metadata={"error": str(exc)},
        )
